#!/usr/bin/env python3
"""Smoke test for the Python `python-pptx` package (import name: pptx).

Fully self-contained: no input files, no network, no packages beyond
python-pptx (and its implied deps). Exercises the core API surface and exits 0
only if every check passes, so it can be used as a pass/fail library validator:

    python3 python-pptx.py

Install: pip install python-pptx   (import name: pptx)

This is the Python counterpart of data.table.R — same contract: a hard
not-installed guard (exit 1), a per-test harness that isolates failures, and a
PASS/exit-0 vs FAIL/exit-1 summary.
"""
import os
import sys

# This file is named after the package it tests, so it sits next to (and would
# shadow) the real top-level module. Drop this script's own directory from the
# import path before importing the package under test.
_here = os.path.dirname(os.path.abspath(__file__))
sys.path = [p for p in sys.path if p not in ("", ".") and os.path.abspath(p) != _here]

try:
    import pptx
except ImportError:
    print("FAIL: package 'python-pptx' is not installed")
    sys.exit(1)

import tempfile

from pptx import Presentation
from pptx.util import Inches


def _version(mod, dist):
    """Best-effort version string: module.__version__, else installed metadata."""
    v = getattr(mod, "__version__", None)
    if v:
        return v
    try:
        import importlib.metadata as m

        return m.version(dist)
    except Exception:
        return "unknown"


print(f"python-pptx version: {_version(pptx, 'python-pptx')}")

failures = 0


def run_test(name, fn):
    """Run one check; a raised exception is a failure, not a crash."""
    global failures
    try:
        fn()
    except Exception as e:  # noqa: BLE001 - any failure is a test failure
        failures += 1
        print(f"  FAIL {name}: {e}")
    else:
        print(f"  ok   {name}")


def test_new_presentation_is_empty():
    prs = Presentation()
    assert len(prs.slides) == 0
    # The default template ships the standard 11 built-in slide layouts.
    assert len(prs.slide_layouts) == 11


def test_add_slide_and_set_title():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    assert len(prs.slides) == 1
    # Layout 0 (title slide) seeds two placeholders: title + subtitle.
    assert len(slide.shapes) == 2
    slide.shapes.title.text = "Smoke Test Title"
    assert slide.shapes.title.text == "Smoke Test Title"


def test_add_textbox_grows_shape_count():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Deck"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "body text"
    assert len(slide.shapes) == 3
    assert box.text_frame.text == "body text"


def test_save_reopen_roundtrip_tempfile():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Persisted Title"
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    try:
        prs.save(path)
        assert os.path.getsize(path) > 0
        reopened = Presentation(path)
        assert len(reopened.slides) == 1
        assert reopened.slides[0].shapes.title.text == "Persisted Title"
    finally:
        os.remove(path)


run_test("new presentation is empty", test_new_presentation_is_empty)
run_test("add slide + set title text", test_add_slide_and_set_title)
run_test("add textbox grows shape count", test_add_textbox_grows_shape_count)
run_test("save/reopen roundtrip tempfile", test_save_reopen_roundtrip_tempfile)

if failures > 0:
    print(f"FAIL: {failures} test(s) failed")
    sys.exit(1)
print("PASS: all python-pptx smoke tests passed")
